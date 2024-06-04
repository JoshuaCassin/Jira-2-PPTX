import csv
from pptx.util import Pt, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from lxml.etree import Element as OxmlElement
from lxml.etree import SubElement
from pptx.dml.color import RGBColor
from datetime import datetime
from pptx.enum.text import PP_ALIGN




namespaces = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}

#add a date and time stamp to each slide
def add_timestamp_to_title_slide(prs):
    slide = prs.slides[0]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Define text box dimensions
    width = Inches(5)
    height = Inches(0.5)
    right = slide_width - width - Inches(0.5)  # 0.5 inch from right
    # left = (slide_width - width) / 2
    top = slide_height - height - Inches(0.3)  # 0.5 inch from bottom

    # Add textbox shape and set properties
    textbox = slide.shapes.add_textbox(right, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "** Current as of " + datetime.utcnow().strftime('%m/%d/%Y %H:%M GMT 0') + " **"
    
    # Style settings
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Set to white
        paragraph.alignment = PP_ALIGN.CENTER

# Define column widths
col_widths = {
    'Issue Key': Inches(.85),
    'SUMMARY': Inches(4.4),
    'FixVersion': Inches(1),
    '': Inches(0.03),
    'CE': Inches(0.4),
    'EE': Inches(0.4)
}

def _set_cell_border(cell, border_color="transparent", border_width='12700'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Define border tags
    border_tags = ['lnL', 'lnR', 'lnT', 'lnB']

    # Check and remove existing borders
    for tag in border_tags:
        existing_border = tcPr.find(".//a:" + tag, namespaces=namespaces)
        if existing_border is not None:
            tcPr.remove(existing_border)

    # Add new borders with correct namespace
    for lines in border_tags:
        ln = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}" + lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
        solidFill = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
        
        if border_color == "transparent":
            noFill = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}noFill")
            solidFill.append(noFill)
        else:
            srgbClr = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr", val=border_color)
            solidFill.append(srgbClr)
        
        prstDash = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash", val='solid')
        round_ = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}round")
        headEnd = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd", type='none', w='med', len='med')
        tailEnd = OxmlElement("{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd", type='none', w='med', len='med')
        
        # Append elements
        ln.append(solidFill)
        ln.append(prstDash)
        ln.append(round_)
        ln.append(headEnd)
        ln.append(tailEnd)

        # Append the border to the cell properties
        tcPr.append(ln)


def colorize_ticks_and_clear_shading(cell):
    # Remove cell shading by setting to transparent
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(31, 41, 56)  # This sets the color to white
    cell.fill.solid()  # This sets the fill type to solid
    cell.fill.opacity = 0.0  # This will set shading to transparent


    # Set borders to transparent (which is the same as no fill)
    _set_cell_border(cell)

    # Apply color to text based on content
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(12)  # <-- Set font size
        for run in paragraph.runs:
            if cell.text == '-':
                run.font.color.rgb = RGBColor(255, 255, 255)  # Green color
            else:
                run.font.color.rgb = RGBColor(0, 128, 0)  # White color



def modify_pptx(file, fix_version_ce, fix_version_ee, display_fields, rows_per_slide, output_path):
    prs = Presentation(file)

    if "CE" not in display_fields:
        display_fields.append("CE")
    if "EE" not in display_fields:
        display_fields.append("EE")

    # Handle the first slide
    slide = prs.slides[0]
    ce_number = fix_version_ce[2:].replace('-', '')
    ee_number = fix_version_ee[2:].replace('-', '')
    title_text = f"Fixversion {ce_number}\nRelease Summary"
    title_placeholder = slide.shapes.title
    if title_placeholder:
        title_placeholder.text = title_text
        for paragraph in title_placeholder.text_frame.paragraphs:
            paragraph.font.size = Pt(38)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    original_width = title_placeholder.width
    original_height = title_placeholder.height

    title_placeholder.left = Inches(0.5)
    title_placeholder.top = Inches(2)
    title_placeholder.text_frame.text_anchor = MSO_ANCHOR.MIDDLE

    for paragraph in title_placeholder.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT

    title_placeholder.width = original_width
    title_placeholder.height = original_height

    add_timestamp_to_title_slide(prs)

    # Add subsequent slides with table data from CSV
    slide_layout = prs.slide_layouts[1]
    with open('output.csv', 'r') as csvfile:
        reader = csv.reader(csvfile)
        headers = next(reader)
        data_rows = [row for row in reader if row]

        for i in range(0, len(data_rows), rows_per_slide):
            slide = prs.slides.add_slide(slide_layout)
            rows_to_add = data_rows[i:i + rows_per_slide]

            # Assuming placeholder[1] is where you want the table
            table_placeholder = slide.shapes[1]

            # Checking if table can be added in the placeholder
            if not hasattr(table_placeholder, 'insert_table'):
                continue

            rows = len(rows_to_add) + 1
            cols = len(display_fields)
            table_shape = table_placeholder.insert_table(rows, cols)
            table = table_shape.table

            # Setting column widths
            for col, field in enumerate(display_fields):
                table.columns[col].width = col_widths.get(field, Inches(1))  # default width is 1 inch

            # Formatting for headers
            for col, field in enumerate(display_fields):
                cell = table.cell(0, col)
                cell.text = field
                colorize_ticks_and_clear_shading(cell)   # <-- Apply formatting to header cells
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.size = Pt(12)

            # Adding data to the table
            for row_idx, data_row in enumerate(rows_to_add, start=1):
                for col, field in enumerate(display_fields):
                    field_idx = headers.index(field)
                    cell_content = data_row[field_idx]
                    cell = table.cell(row_idx, col)
                    cell.text = cell_content
                    
                    colorize_ticks_and_clear_shading(cell)

         # The loop to remove all borders from all tables across the presentation
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        colorize_ticks_and_clear_shading(cell)    
        
   

    prs.save(output_path)
